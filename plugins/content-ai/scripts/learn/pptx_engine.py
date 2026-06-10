"""
Reusable PPTX generation engine.
Creates branded PowerPoint presentations from structured data.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from themes.anthropic import THEME as DEFAULT_THEME


# Widescreen slide dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class PptxEngine:
    """Generates branded PowerPoint presentations from structured content."""

    def __init__(self, title, subtitle="", footer_text="", theme=None):
        self.theme = theme or DEFAULT_THEME
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.slide_num = 0
        self.footer_text = footer_text or title

        # Create title slide
        self._make_title_slide(title, subtitle)

    def save(self, output_path):
        self.prs.save(output_path)
        return len(self.prs.slides)

    # --- Low-level helpers ---

    def _blank_slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _add_bg(self, slide, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def _add_accent_bar(self, slide, top, color=None, height=Inches(0.06)):
        color = color or self.theme["accent"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), top, SLIDE_W, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def _add_side_stripe(self, slide, width=Inches(0.35)):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), width, SLIDE_H
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme["stripe_color"]
        shape.line.fill.background()

    def _add_footer(self, slide):
        txBox = slide.shapes.add_textbox(
            Inches(8.5), Inches(7.0), Inches(4.5), Inches(0.4)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = self.footer_text
        p.font.size = Pt(10)
        p.font.color.rgb = self.theme["footer_text"]
        p.alignment = PP_ALIGN.RIGHT

    def _add_slide_number(self, slide):
        txBox = slide.shapes.add_textbox(
            Inches(0.6), Inches(7.0), Inches(1.0), Inches(0.4)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(self.slide_num)
        p.font.size = Pt(10)
        p.font.color.rgb = self.theme["footer_text"]
        p.alignment = PP_ALIGN.LEFT

    def _set_cell(self, cell, text, font_size=Pt(13), bold=False, color=None,
                  alignment=PP_ALIGN.LEFT):
        color = color or self.theme["body_text"]
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        cell.text_frame.word_wrap = True
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)

    # --- Slide builders ---

    def _make_title_slide(self, title, subtitle):
        self.slide_num += 1
        slide = self._blank_slide()
        self._add_bg(slide, self.theme["title_bg"])
        self._add_accent_bar(slide, Inches(3.4), height=Inches(0.06))

        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.theme["title_text"]
        p.alignment = PP_ALIGN.LEFT

        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(3.7), Inches(10), Inches(1.2))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = self.theme["subtitle_text"]
            p2.alignment = PP_ALIGN.LEFT

        self._add_footer(slide)

    def add_section_slide(self, section_num, title):
        self.slide_num += 1
        slide = self._blank_slide()
        self._add_bg(slide, self.theme["section_bg"])
        self._add_accent_bar(slide, Inches(4.0), height=Inches(0.06))

        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(3), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"SECTION {section_num:02d}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.theme["accent"]
        p.alignment = PP_ALIGN.LEFT

        txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(11), Inches(1.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(40)
        p2.font.bold = True
        p2.font.color.rgb = self.theme["title_text"]
        p2.alignment = PP_ALIGN.LEFT

        self._add_footer(slide)

    def add_content_slide(self, title, bullets, sub_bullets=None, font_size=Pt(16)):
        self.slide_num += 1
        slide = self._blank_slide()
        self._add_bg(slide, self.theme["content_bg"])
        self._add_side_stripe(slide)
        self._add_accent_bar(slide, Inches(1.15), color=self.theme["heading_text"],
                             height=Inches(0.03))

        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(12), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme["heading_text"]
        p.alignment = PP_ALIGN.LEFT

        txBox2 = slide.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.8), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        if sub_bullets is None:
            sub_bullets = {}

        for i, bullet in enumerate(bullets):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.text = bullet
            p.font.size = font_size
            p.font.color.rgb = self.theme["body_text"]
            p.space_after = Pt(6)
            p.level = 0
            if i in sub_bullets:
                for sb in sub_bullets[i]:
                    sp = tf2.add_paragraph()
                    sp.text = f"    {sb}"
                    sp.font.size = Pt(14)
                    sp.font.color.rgb = self.theme["muted_text"]
                    sp.space_after = Pt(3)
                    sp.level = 1

        self._add_slide_number(slide)
        self._add_footer(slide)

    def add_table_slide(self, title, headers, rows, col_widths=None,
                        font_size=Pt(12), header_font_size=Pt(13)):
        self.slide_num += 1
        slide = self._blank_slide()
        self._add_bg(slide, self.theme["content_bg"])
        self._add_side_stripe(slide)
        self._add_accent_bar(slide, Inches(1.15), color=self.theme["heading_text"],
                             height=Inches(0.03))

        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(12), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme["heading_text"]

        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl_width = Inches(12.0)
        tbl_left = Inches(0.7)
        tbl_top = Inches(1.4)
        row_height = Inches(0.45)
        tbl_height = row_height * n_rows

        table_shape = slide.shapes.add_table(
            n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height
        )
        table = table_shape.table

        if col_widths:
            for i, w in enumerate(col_widths):
                table.columns[i].width = Inches(w) if isinstance(w, (int, float)) else w
        else:
            default_w = int(tbl_width / n_cols)
            for i in range(n_cols):
                table.columns[i].width = default_w

        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            self._set_cell(cell, h, font_size=header_font_size, bold=True,
                           color=self.theme["table_header_text"])
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme["table_header_bg"]

        for r_idx, row in enumerate(rows):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                self._set_cell(cell, val, font_size=font_size)
                cell.fill.solid()
                if r_idx % 2 == 0:
                    cell.fill.fore_color.rgb = self.theme["table_row_even"]
                else:
                    cell.fill.fore_color.rgb = self.theme["table_row_odd"]

        self._add_slide_number(slide)
        self._add_footer(slide)

    def add_two_column_slide(self, title, left_title, left_items, right_title, right_items):
        self.slide_num += 1
        slide = self._blank_slide()
        self._add_bg(slide, self.theme["content_bg"])
        self._add_side_stripe(slide)
        self._add_accent_bar(slide, Inches(1.15), color=self.theme["heading_text"],
                             height=Inches(0.03))

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(12), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.theme["heading_text"]

        # Left column header
        left_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.6)
        )
        left_header.fill.solid()
        left_header.fill.fore_color.rgb = self.theme["accent"]
        left_header.line.fill.background()
        ltf = left_header.text_frame
        ltf.paragraphs[0].text = left_title
        ltf.paragraphs[0].font.size = Pt(18)
        ltf.paragraphs[0].font.bold = True
        ltf.paragraphs[0].font.color.rgb = self.theme["title_text"]
        ltf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ltf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Right column header
        right_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.4), Inches(5.8), Inches(0.6)
        )
        right_header.fill.solid()
        right_header.fill.fore_color.rgb = RGBColor(0x4A, 0x90, 0xD9)  # accent blue
        right_header.line.fill.background()
        rtf = right_header.text_frame
        rtf.paragraphs[0].text = right_title
        rtf.paragraphs[0].font.size = Pt(18)
        rtf.paragraphs[0].font.bold = True
        rtf.paragraphs[0].font.color.rgb = self.theme["title_text"]
        rtf.paragraphs[0].alignment = PP_ALIGN.CENTER
        rtf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Left bullets
        ltxBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5))
        ltf2 = ltxBox.text_frame
        ltf2.word_wrap = True
        for i, item in enumerate(left_items):
            p = ltf2.paragraphs[0] if i == 0 else ltf2.add_paragraph()
            p.text = item
            p.font.size = Pt(15)
            p.font.color.rgb = self.theme["body_text"]
            p.space_after = Pt(6)

        # Right bullets
        rtxBox = slide.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.5), Inches(4.5))
        rtf2 = rtxBox.text_frame
        rtf2.word_wrap = True
        for i, item in enumerate(right_items):
            p = rtf2.paragraphs[0] if i == 0 else rtf2.add_paragraph()
            p.text = item
            p.font.size = Pt(15)
            p.font.color.rgb = self.theme["body_text"]
            p.space_after = Pt(6)

        self._add_slide_number(slide)
        self._add_footer(slide)

    def add_closing_slide(self, title="Thank You", subtitle="", resources=None):
        """Final slide with dark background."""
        self.slide_num += 1
        slide = self._blank_slide()
        self._add_bg(slide, self.theme["title_bg"])
        self._add_accent_bar(slide, Inches(3.6), height=Inches(0.04))

        txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.theme["title_text"]

        if resources:
            txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(2.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for i, r in enumerate(resources):
                p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                p.text = r
                p.font.size = Pt(18)
                p.font.color.rgb = self.theme["subtitle_text"]
                p.space_after = Pt(8)

        if subtitle:
            txBox3 = slide.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(11), Inches(1.5))
            tf3 = txBox3.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = subtitle
            p3.font.size = Pt(40)
            p3.font.bold = True
            p3.font.color.rgb = self.theme["accent"]

        self._add_footer(slide)
